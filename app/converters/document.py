from __future__ import annotations

import csv
import json
import logging
import os
import shutil
import subprocess
import tempfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any, Dict, List, Optional

from .base import BaseConverter

logger = logging.getLogger(__name__)

try:
    from PyPDF2 import PdfReader, PdfWriter
    _PYPDF2 = True
except ImportError:
    _PYPDF2 = False

try:
    import docx as _docx_module
    _DOCX = True
except ImportError:
    _DOCX = False

try:
    from pdf2docx import Converter as Pdf2DocxConverter
    _PDF2DOCX = True
except ImportError:
    _PDF2DOCX = False

try:
    import markdown as _md_module
    _MARKDOWN = True
except ImportError:
    _MARKDOWN = False

try:
    import html2text as _h2t_module
    _HTML2TEXT = True
except ImportError:
    _HTML2TEXT = False

try:
    from bs4 import BeautifulSoup
    _BS4 = True
except ImportError:
    _BS4 = False

try:
    import openpyxl
    _OPENPYXL = True
except ImportError:
    _OPENPYXL = False

try:
    from pptx import Presentation as _PptxPresentation
    _PPTX = True
except ImportError:
    _PPTX = False

try:
    from weasyprint import HTML as WeasyHTML
    _WEASYPRINT = True
except Exception:
    _WEASYPRINT = False

try:
    from reportlab.lib.pagesizes import A4, letter
    from reportlab.lib.units import mm
    from reportlab.pdfgen import canvas as rl_canvas
    import textwrap as _textwrap
    _REPORTLAB = True
except ImportError:
    _REPORTLAB = False


def _check_bin(name: str) -> bool:
    return shutil.which(name) is not None


_PANDOC = _check_bin("pandoc")
_LIBREOFFICE = _check_bin("libreoffice") or _check_bin("soffice")


class DocumentConverter(BaseConverter):

    def get_supported_input_formats(self) -> List[str]:
        base = [
            "pdf", "docx", "doc", "odt", "txt", "rtf",
            "html", "htm", "md", "csv", "json", "xml",
        ]
        if _OPENPYXL:
            base += ["xlsx", "xls"]
        if _PPTX or _LIBREOFFICE:
            base += ["pptx"]
        if _PANDOC:
            base += ["epub", "tex", "org", "rst", "adoc", "ods"]
        return sorted(set(base))

    def get_supported_output_formats(self) -> List[str]:
        base = [
            "pdf", "docx", "txt", "html", "md", "csv", "json", "xml",
        ]
        if _PANDOC:
            base += ["epub", "tex", "org", "rst", "odt", "rtf"]
        if _OPENPYXL:
            base += ["xlsx"]
        return sorted(set(base))

    def convert(self, input_path: str, output_path: str, **options: Any) -> bool:
        src = self.src_ext(input_path)
        dst = self.dst_ext(output_path)
        try:

            if dst == "pdf":
                return self._to_pdf(input_path, output_path, src, options)
            if src == "pdf":
                return self._from_pdf(input_path, output_path, dst, options)

            if src in ("xlsx", "xls") and dst in ("csv", "json", "txt"):
                return self._xlsx_to_flat(input_path, output_path, dst)
            if src == "csv" and dst in ("xlsx",):
                return self._csv_to_xlsx(input_path, output_path)

            if src == "pptx" and dst in ("txt", "html", "pdf"):
                return self._pptx_convert(input_path, output_path, dst, options)

            if src in ("docx", "doc", "odt", "rtf") and dst in ("docx", "doc", "odt", "rtf", "txt"):
                return self._office_convert(input_path, output_path, src, dst)


            if {src, dst} & {"txt", "md", "html", "htm"}:
                return self._text_convert(input_path, output_path, src, dst)

            if src in ("csv", "json", "xml") or dst in ("csv", "json", "xml"):
                return self._structured_convert(input_path, output_path, src, dst)


            if _PANDOC:
                return self._pandoc(input_path, output_path, options)

            logger.error("No conversion path for %s → %s", src, dst)
            return False

        except Exception:
            logger.exception("Document conversion failed: %s → %s", input_path, output_path)
            return False

    def _to_pdf(self, inp: str, out: str, src: str, opts: dict) -> bool:
        if _LIBREOFFICE and src in ("docx", "doc", "odt", "rtf", "xlsx", "pptx", "ods", "html", "htm"):
            ok = self._libreoffice_to_pdf(inp, out)
            if ok:
                return True
            logger.warning("LibreOffice failed, trying next method")

        if src in ("md", "markdown", "txt", "html", "htm"):
            return self._html_to_pdf_weasyprint(inp, out, src, opts)

        if _PANDOC:
            return self._pandoc(inp, out, opts)

        if src in ("txt", "md") and _REPORTLAB:
            return self._txt_to_pdf_reportlab(inp, out, opts)

        logger.error("No PDF output path available for source format %s", src)
        return False

    def _libreoffice_to_pdf(self, inp: str, out: str) -> bool:
        out_dir = str(Path(out).parent)
        bin_name = "libreoffice" if shutil.which("libreoffice") else "soffice"
        cmd = [bin_name, "--headless", "--convert-to", "pdf",
               "--outdir", out_dir, inp]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if result.returncode != 0:
            logger.error("LibreOffice error: %s", result.stderr)
            return False
        expected = Path(out_dir) / (Path(inp).stem + ".pdf")
        if expected.exists() and str(expected) != out:
            shutil.move(str(expected), out)
        return Path(out).exists()

    def _html_to_pdf_weasyprint(self, inp: str, out: str, src: str, opts: dict) -> bool:
        if not _WEASYPRINT:
            if src in ("txt", "md") and _REPORTLAB:
                return self._txt_to_pdf_reportlab(inp, out, opts)
            logger.error("WeasyPrint not installed; cannot convert %s → pdf", src)
            return False

        with open(inp, "r", encoding="utf-8", errors="replace") as fh:
            content = fh.read()

        if src in ("md", "markdown"):
            if _MARKDOWN:
                html = _md_module.markdown(content, extensions=["tables", "fenced_code"])
            else:
                html = f"<pre>{content}</pre>"
        elif src == "txt":
            html = f"<html><body><pre style='white-space:pre-wrap'>{content}</pre></body></html>"
        else:
            html = content

        if not html.strip().lower().startswith("<!doctype") and "<html" not in html.lower():
            css = self._default_css(opts)
            html = f"<!doctype html><html><head><meta charset='utf-8'><style>{css}</style></head><body>{html}</body></html>"

        WeasyHTML(string=html, base_url=str(Path(inp).parent)).write_pdf(out)
        return Path(out).exists()

    @staticmethod
    def _default_css(opts: dict) -> str:
        paper = opts.get("paper_size", "A4").lower()
        margin = opts.get("margins", {})
        mt = margin.get("top", 20)
        mr = margin.get("right", 20)
        mb = margin.get("bottom", 20)
        ml = margin.get("left", 20)
        font = opts.get("font", "Arial, sans-serif")
        font_size = opts.get("font_size", 12)
        return (
            f"@page {{ size: {paper}; margin: {mt}mm {mr}mm {mb}mm {ml}mm; }}"
            f"body {{ font-family: {font}; font-size: {font_size}pt; line-height: 1.5; }}"
            "h1,h2,h3 { margin-top: 1em; } pre { white-space: pre-wrap; }"
        )

    def _txt_to_pdf_reportlab(self, inp: str, out: str, opts: dict) -> bool:
        if not _REPORTLAB:
            return False
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
        import textwrap
        with open(inp, "r", encoding="utf-8", errors="replace") as fh:
            text = fh.read()
        c = canvas.Canvas(out, pagesize=A4)
        w, h = A4
        x, y = 20 * mm, h - 20 * mm
        c.setFont("Helvetica", 11)
        for para in text.splitlines():
            for line in textwrap.wrap(para, 95) or [""]:
                if y < 20 * mm:
                    c.showPage()
                    y = h - 20 * mm
                    c.setFont("Helvetica", 11)
                c.drawString(x, y, line)
                y -= 14
            y -= 4
        c.save()
        return Path(out).exists()

    def _from_pdf(self, inp: str, out: str, dst: str, opts: dict) -> bool:
        if dst == "docx":
            return self._pdf_to_docx(inp, out)
        if dst in ("txt", "md", "html"):
            return self._pdf_to_text(inp, out, dst, opts)
        if _PANDOC:
            return self._pandoc(inp, out, opts)
        logger.error("No path from PDF to %s", dst)
        return False

    def _pdf_to_docx(self, inp: str, out: str) -> bool:
        if not _PDF2DOCX:
            if _LIBREOFFICE:
                return self._libreoffice_to_format(inp, out, "docx")
            logger.error("pdf2docx not installed; PDF→DOCX unavailable")
            return False
        try:
            cv = Pdf2DocxConverter(inp)
            cv.convert(out)
            cv.close()
            return Path(out).exists()
        except Exception:
            logger.exception("pdf2docx failed")
            return False

    def _libreoffice_to_format(self, inp: str, out: str, ext: str) -> bool:
        out_dir = str(Path(out).parent)
        bin_name = "libreoffice" if shutil.which("libreoffice") else "soffice"
        cmd = [bin_name, "--headless", "--convert-to", ext, "--outdir", out_dir, inp]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        expected = Path(out_dir) / (Path(inp).stem + f".{ext}")
        if expected.exists() and str(expected) != out:
            shutil.move(str(expected), out)
        return Path(out).exists()

    def _pdf_to_text(self, inp: str, out: str, dst: str, opts: dict) -> bool:
        if not _PYPDF2:
            logger.error("PyPDF2 not installed")
            return False
        reader = PdfReader(inp)
        if reader.is_encrypted:
            pw = opts.get("password", "")
            if not reader.decrypt(pw):
                logger.error("PDF decryption failed")
                return False
        pages = [p.extract_text() or "" for p in reader.pages]
        text = "\n\n".join(pages)
        if dst == "html":
            import html as _html_lib
            body = "<br>".join(_html_lib.escape(ln) for ln in text.splitlines())
            content = f"<!doctype html><html><body><p>{body}</p></body></html>"
        elif dst == "md":
            content = "# Extracted from PDF\n\n" + text
        else:
            content = text
        with open(out, "w", encoding="utf-8") as fh:
            fh.write(content)
        return True

    def _office_convert(self, inp: str, out: str, src: str, dst: str) -> bool:
        if _LIBREOFFICE:
            return self._libreoffice_to_format(inp, out, dst)
        if _PANDOC:
            return self._pandoc(inp, out, {})
        if src == "docx" and dst == "txt" and _DOCX:
            doc = _docx_module.Document(inp)
            with open(out, "w", encoding="utf-8") as fh:
                fh.write("\n".join(p.text for p in doc.paragraphs))
            return True
        logger.error("No Office conversion path for %s → %s", src, dst)
        return False

    def _text_convert(self, inp: str, out: str, src: str, dst: str) -> bool:
        with open(inp, "r", encoding="utf-8", errors="replace") as fh:
            content = fh.read()

        if src in ("md", "markdown") and dst in ("html", "htm"):
            result = _md_module.markdown(content, extensions=["tables", "fenced_code"]) if _MARKDOWN else f"<pre>{content}</pre>"
        elif src in ("html", "htm") and dst in ("md", "markdown"):
            if _HTML2TEXT:
                h = _h2t_module.HTML2Text()
                h.ignore_links = False
                result = h.handle(content)
            else:
                result = BeautifulSoup(content, "html.parser").get_text() if _BS4 else content
        elif dst == "txt":
            if src in ("html", "htm"):
                result = BeautifulSoup(content, "html.parser").get_text() if _BS4 else content
            else:
                result = content.replace("#", "").replace("*", "").replace("_", "").replace("`", "")
        elif src == "txt" and dst in ("html", "htm"):
            result = f"<!doctype html><html><body><pre style='white-space:pre-wrap'>{content}</pre></body></html>"
        elif src == "txt" and dst in ("md", "markdown"):
            lines = content.splitlines()
            parts = [f"# {lines[0]}\n"] if lines else []
            parts += [ln + "\n" for ln in lines[1:]]
            result = "\n".join(parts)
        else:
            if _PANDOC:
                return self._pandoc(inp, out, {})
            result = content

        with open(out, "w", encoding="utf-8") as fh:
            fh.write(result)
        return True

    def _structured_convert(self, inp: str, out: str, src: str, dst: str) -> bool:
        if src == "csv" and dst == "json":
            return self._csv_to_json(inp, out)
        if src == "json" and dst == "csv":
            return self._json_to_csv(inp, out)
        if src == "xml" and dst == "json":
            return self._xml_to_json(inp, out)
        if src == "json" and dst == "xml":
            return self._json_to_xml(inp, out)
        if src == "csv" and dst == "xml":
            tmp = tempfile.mktemp(suffix=".json")
            try:
                return self._csv_to_json(inp, tmp) and self._json_to_xml(tmp, out)
            finally:
                Path(tmp).unlink(missing_ok=True)
        if src == "xml" and dst == "csv":
            tmp = tempfile.mktemp(suffix=".json")
            try:
                return self._xml_to_json(inp, tmp) and self._json_to_csv(tmp, out)
            finally:
                Path(tmp).unlink(missing_ok=True)
        logger.error("Unsupported structured conversion: %s → %s", src, dst)
        return False

    @staticmethod
    def _csv_to_json(inp: str, out: str) -> bool:
        rows: list = []
        with open(inp, newline="", encoding="utf-8-sig") as fh:
            rows = list(csv.DictReader(fh))
        with open(out, "w", encoding="utf-8") as fh:
            json.dump(rows, fh, indent=2, ensure_ascii=False)
        return True

    @staticmethod
    def _json_to_csv(inp: str, out: str) -> bool:
        with open(inp, encoding="utf-8") as fh:
            data = json.load(fh)
        if not isinstance(data, list):
            data = [data]
        if not data:
            Path(out).write_text("")
            return True
        fields = list(dict.fromkeys(k for row in data for k in (row.keys() if isinstance(row, dict) else {})))
        with open(out, "w", newline="", encoding="utf-8") as fh:
            w = csv.DictWriter(fh, fieldnames=fields, extrasaction="ignore")
            w.writeheader()
            w.writerows(data)
        return True

    @staticmethod
    def _xml_to_json(inp: str, out: str) -> bool:
        def _parse(elem: ET.Element) -> Any:
            d: dict = {}
            if elem.attrib:
                d["@attr"] = elem.attrib
            children = list(elem)
            if not children:
                d["#text"] = (elem.text or "").strip()
            else:
                for child in children:
                    tag, val = child.tag, _parse(child)
                    if tag in d:
                        if not isinstance(d[tag], list):
                            d[tag] = [d[tag]]
                        d[tag].append(val)
                    else:
                        d[tag] = val
            return d
        root = ET.parse(inp).getroot()
        with open(out, "w", encoding="utf-8") as fh:
            json.dump({root.tag: _parse(root)}, fh, indent=2, ensure_ascii=False)
        return True

    @staticmethod
    def _json_to_xml(inp: str, out: str) -> bool:
        def _build(data: Any, tag: str) -> ET.Element:
            elem = ET.Element(tag)
            if isinstance(data, dict):
                for k, v in data.items():
                    if k == "@attr" and isinstance(v, dict):
                        for ak, av in v.items():
                            elem.set(ak, str(av))
                    elif k == "#text":
                        elem.text = str(v)
                    else:
                        elem.append(_build(v, k))
            elif isinstance(data, list):
                for item in data:
                    elem.append(_build(item, "item"))
            else:
                elem.text = str(data)
            return elem

        with open(inp, encoding="utf-8") as fh:
            data = json.load(fh)
        if isinstance(data, dict) and len(data) == 1:
            root_key = next(iter(data))
            root = _build(data[root_key], root_key)
        else:
            root = _build(data, "root")
        ET.indent(root)
        ET.ElementTree(root).write(out, encoding="utf-8", xml_declaration=True)
        return True

    def _xlsx_to_flat(self, inp: str, out: str, dst: str) -> bool:
        if not _OPENPYXL:
            if _LIBREOFFICE:
                return self._libreoffice_to_format(inp, out, dst)
            logger.error("openpyxl not installed")
            return False
        wb = openpyxl.load_workbook(inp, data_only=True)
        ws = wb.active
        if ws is None:
            logger.error("openpyxl: workbook has no active sheet")
            return False
        rows = [[str(c.value) if c.value is not None else "" for c in row] for row in ws.iter_rows()]
        if dst == "csv":
            with open(out, "w", newline="", encoding="utf-8") as fh:
                csv.writer(fh).writerows(rows)
        elif dst == "json":
            if rows:
                headers = rows[0]
                data = [dict(zip(headers, r)) for r in rows[1:]]
            else:
                data = []
            with open(out, "w", encoding="utf-8") as fh:
                json.dump(data, fh, indent=2, ensure_ascii=False)
        elif dst == "txt":
            with open(out, "w", encoding="utf-8") as fh:
                for r in rows:
                    fh.write("\t".join(r) + "\n")
        return True

    def _csv_to_xlsx(self, inp: str, out: str) -> bool:
        if not _OPENPYXL:
            logger.error("openpyxl not installed")
            return False
        wb = openpyxl.Workbook()
        ws = wb.active
        if ws is None:
            logger.error("openpyxl: new workbook has no active sheet")
            return False
        with open(inp, newline="", encoding="utf-8-sig") as fh:
            for row in csv.reader(fh):
                ws.append(row)
        wb.save(out)
        return True

    def _pptx_convert(self, inp: str, out: str, dst: str, opts: dict) -> bool:
        if _LIBREOFFICE and dst == "pdf":
            return self._libreoffice_to_pdf(inp, out)
        if _PPTX and dst == "txt":
            prs = _PptxPresentation(inp)
            lines: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    txt = getattr(shape, "text", "")
                    if txt:
                        lines.append(txt)
            with open(out, "w", encoding="utf-8") as fh:
                fh.write("\n".join(lines))
            return True
        if _PANDOC:
            return self._pandoc(inp, out, opts)
        return False

    def _pandoc(self, inp: str, out: str, opts: dict) -> bool:
        if not _PANDOC:
            return False
        cmd = ["pandoc", inp, "-o", out, "--standalone"]
        if opts.get("toc"):
            cmd.append("--toc")
        if opts.get("paper_size"):
            cmd += ["-V", f"papersize={opts['paper_size']}"]
        if opts.get("font"):
            cmd += ["-V", f"mainfont={opts['font']}"]
        if opts.get("font_size"):
            cmd += ["-V", f"fontsize={opts['font_size']}pt"]
        margins = opts.get("margins", {})
        for side in ("top", "right", "bottom", "left"):
            if side in margins:
                cmd += ["-V", f"margin-{side}={margins[side]}mm"]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if result.returncode != 0:
            logger.error("Pandoc error: %s", result.stderr)
            return False
        return Path(out).exists()
